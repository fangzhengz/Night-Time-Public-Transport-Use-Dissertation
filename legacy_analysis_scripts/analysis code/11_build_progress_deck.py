# -*- coding: utf-8 -*-
"""Build RQ1 progress-update deck, matching the existing night_transport deck style."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image
import os

OUT = r"D:\SDS2025_workspace\CASA_FYP\FYP\night_transport_RQ1_progress.pptx"
IMG = r"D:\SDS2025_workspace\CASA_FYP\FYP\outputs"

# palette (match reference: purple accent 500778)
DARK   = RGBColor(0x1F, 0x16, 0x40)   # title / closing bg
ACCENT = RGBColor(0x50, 0x07, 0x78)   # deep purple
ACC2   = RGBColor(0x8E, 0x6B, 0xB8)   # light purple
BODY   = RGBColor(0x33, 0x33, 0x33)
MUTED  = RGBColor(0x66, 0x66, 0x66)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
TINT   = RGBColor(0xF4, 0xF0, 0xF8)   # light purple card
GREEN  = RGBColor(0x2F, 0x6B, 0x4F)
RED    = RGBColor(0x9A, 0x3D, 0x3D)

prs = Presentation()
prs.slide_width = Inches(10); prs.slide_height = Inches(5.625)
BLANK = prs.slide_layouts[6]
EMU = 914400


def slide(bg=WHITE):
    s = prs.slides.add_slide(BLANK)
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    r.fill.solid(); r.fill.fore_color.rgb = bg; r.line.fill.background()
    r.shadow.inherit = False
    s.shapes._spTree.remove(r._element); s.shapes._spTree.insert(2, r._element)
    return s


def tb(s, l, t, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, sp_after=4, line=1.0):
    """runs: list of paragraphs; each paragraph = list of (text,size,color,bold,italic)."""
    box = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = box.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = Inches(0.05); tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02); tf.margin_bottom = Inches(0.02)
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.space_after = Pt(sp_after); p.space_before = Pt(0)
        p.line_spacing = line
        for (text, size, color, bold, italic) in para:
            r = p.add_run(); r.text = text
            r.font.name = "Calibri"; r.font.size = Pt(size)
            r.font.color.rgb = color; r.font.bold = bold; r.font.italic = italic
    return box


def rect(s, l, t, w, h, fill, line=None):
    r = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    r.fill.solid(); r.fill.fore_color.rgb = fill
    if line is None:
        r.line.fill.background()
    else:
        r.line.color.rgb = line; r.line.width = Pt(1)
    r.shadow.inherit = False
    return r


def accent_bar(s, t):
    b = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.4), Inches(t), Inches(0.08), Inches(0.42))
    b.fill.solid(); b.fill.fore_color.rgb = ACCENT; b.line.fill.background(); b.shadow.inherit = False


def title(s, text, sub=None):
    accent_bar(s, 0.30)
    tb(s, 0.6, 0.22, 9.0, 0.6, [[(text, 23, ACCENT, True, False)]])
    if sub:
        tb(s, 0.6, 0.80, 9.0, 0.4, [[(sub, 12.5, MUTED, False, True)]])


def pic_fit(s, path, l, t, maxw, maxh, caption=None):
    iw, ih = Image.open(path).size
    ar = iw / ih
    w, h = maxw, maxw / ar
    if h > maxh:
        h = maxh; w = maxh * ar
    left = l + (maxw - w) / 2
    s.shapes.add_picture(path, Inches(left), Inches(t), Inches(w), Inches(h))
    if caption:
        tb(s, l, t + h + 0.02, maxw, 0.3, [[(caption, 8.5, MUTED, False, True)]], align=PP_ALIGN.CENTER)
    return h


def P(IMGsub):
    return os.path.join(IMG, IMGsub)


# ---------- Slide 1: title ----------
s = slide(DARK)
tb(s, 0.6, 1.55, 8.8, 0.6, [[("Night-time Public Transport & Spatial Equity in London", 27, WHITE, True, False)]])
tb(s, 0.6, 2.55, 8.8, 0.9, [
    [("RQ1 Progress Update — Clustering Methodology & Tentative Results", 15, ACC2, True, False)],
])
tb(s, 0.6, 3.35, 8.8, 0.5, [[("Exploration process · methodological changes · current results", 12.5, RGBColor(0xCF,0xC4,0xE4), False, True)]])
tb(s, 0.6, 4.55, 8.8, 0.4, [[("Fangzheng  |  UCL CASA × TfL", 12, RGBColor(0xB9,0xAD,0xD6), False, False)]])

# ---------- Slide 2: where we are ----------
s = slide()
title(s, "Where we are — RQ1 in one line", "Group stations & areas by their night-time demand 'fingerprint'")
tb(s, 0.6, 1.25, 9.0, 0.7, [[
    ("Goal: ", 14, ACCENT, True, False),
    ("identify distinct night-time use profiles by clustering stations / LSOAs on their temporal demand curves (20:00-start window), as the basis for the later equity analysis.", 14, BODY, False, False)]], line=1.1)
cards = [
    ("Rail (NUMBAT)", "270 LU stations\nentry / exit, 15-min,\nMON–SUN"),
    ("Bus (BUSTO)", "Full network → LSOA\nboardings / alightings\nWeekday / Sat / Sun"),
    ("Spatial layer", "London LSOA 2021\n+ TfL station\ncoordinates"),
    ("Window", "20:00 → 01:00 (wkday)\n20:00 → 06:00 (wkend)\nlow-volume filtered"),
]
x = 0.5
for h_, body_ in cards:
    rect(s, x, 2.15, 2.2, 1.9, TINT)
    tb(s, x+0.12, 2.28, 1.96, 0.5, [[(h_, 12.5, ACCENT, True, False)]])
    tb(s, x+0.12, 2.78, 1.96, 1.2, [[(body_, 10.5, BODY, False, False)]], line=1.05)
    x += 2.37
tb(s, 0.6, 4.25, 9.0, 0.7, [[
    ("This update focuses on RQ1 only: ", 12, MUTED, True, True),
    ("how the clustering approach evolved and what the current clusters look like.", 12, MUTED, False, True)]], line=1.05)

# ---------- Slide 3: the journey (process flow) ----------
s = slide()
title(s, "The exploration journey", "How the method changed as the data pushed back")
steps = [
    ("1", "Baseline", "GMM on raw 15-min\nentry/exit shares\n(K fixed: 5 / 6)"),
    ("2", "Problem", "K won't separate;\nnear-duplicate clusters,\nsilhouette ≈ 0.07"),
    ("3", "Diagnosis", "Shared night-decline\nenvelope dominates;\nstructure ~continuous"),
    ("4", "Re-design", "Engineer behavioural\nfeatures that isolate\nthe discriminative signal"),
    ("5", "Evidence-based K", "GMM/KMeans/Ward +\nbootstrap stability +\ninterpretability"),
]
x = 0.45; w = 1.74
for i,(n,h_,b_) in enumerate(steps):
    rect(s, x, 1.7, w, 2.4, TINT)
    c = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x+w/2-0.22), Inches(1.85), Inches(0.44), Inches(0.44))
    c.fill.solid(); c.fill.fore_color.rgb = ACCENT; c.line.fill.background(); c.shadow.inherit=False
    tf=c.text_frame; tf.word_wrap=False; p=tf.paragraphs[0]; p.alignment=PP_ALIGN.CENTER
    r=p.add_run(); r.text=n; r.font.name="Calibri"; r.font.size=Pt(16); r.font.bold=True; r.font.color.rgb=WHITE
    tb(s, x+0.05, 2.4, w-0.1, 0.4, [[(h_, 12.5, ACCENT, True, False)]], align=PP_ALIGN.CENTER)
    tb(s, x+0.05, 2.85, w-0.1, 1.2, [[(b_, 10, BODY, False, False)]], align=PP_ALIGN.CENTER, line=1.05)
    if i < 4:
        tb(s, x+w-0.06, 2.55, 0.3, 0.4, [[("→", 18, ACC2, True, False)]])
    x += w + 0.12
tb(s, 0.6, 4.45, 9.0, 0.6, [[
    ("Key shift: ", 12.5, ACCENT, True, False),
    ("stop imposing K; let feature design expose the structure, then choose K from convergent evidence.", 12.5, BODY, False, False)]], line=1.05)

# ---------- Slide 4: problem with baseline ----------
s = slide()
title(s, "Problem with the baseline", "Raw quarter-hour shares could not separate clean clusters")
pic_fit(s, P(r"lu_rail_2000_diag\weekday\weekday_profiles.png"), 0.4, 1.25, 5.0, 4.1,
        caption="Old GMM K=5 (weekday): C0/C1/C3 nearly identical curves")
tb(s, 5.65, 1.35, 4.05, 0.5, [[("What went wrong", 14, ACCENT, True, False)]])
for i,(t_) in enumerate([
    "75% of stations (C0+C1+C3) collapsed into near-duplicate profiles",
    "Centroid cosine distance < 0.018 — clusters differ by level, not shape",
    "Silhouette ≈ 0.07; K=5/6 was hard-coded, not selected by any index",
    "Indices actually favoured very low K — the split was largely arbitrary",
]):
    rect(s, 5.65, 1.9+i*0.78, 0.08, 0.6, ACCENT)
    tb(s, 5.85, 1.88+i*0.78, 3.9, 0.75, [[(t_, 11.5, BODY, False, False)]], line=1.0)

# ---------- Slide 5: diagnosis ----------
s = slide()
title(s, "Diagnosis — why K wouldn't resolve", "A shared signal was drowning the discriminative one")
tb(s, 0.6, 1.25, 9.0, 1.0, [[
    ("Every station's curve carries the same ", 14, BODY, False, False),
    ("night-decline envelope", 14, ACCENT, True, False),
    (" (high at 20:00, fading to ~0). In the raw 80-dim shares this common shape is the largest source of variance, so distance is dominated by overall level — the signals that distinguish station roles (direction, timing, persistence) are buried.", 14, BODY, False, False)]], line=1.12)
boxes = [
    ("Shared envelope = noise for clustering", "The decline shape is common to all → contributes nothing to separation, yet dominates the metric."),
    ("Structure is partly continuous", "Profiles form a gradient, not discrete blobs → silhouette is inherently low; hard K is fragile."),
    ("Old K was imposed, not discovered", "FINAL_K = 5/6 was set manually; the indices alone pointed to far fewer groups."),
]
y = 2.5
for h_,b_ in boxes:
    rect(s, 0.5, y, 9.0, 0.86, TINT)
    tb(s, 0.7, y+0.08, 8.7, 0.4, [[(h_, 12.5, ACCENT, True, False)]])
    tb(s, 0.7, y+0.45, 8.7, 0.4, [[(b_, 11, BODY, False, False)]], line=1.0)
    y += 0.98

# ---------- Slide 6: methodological change - features ----------
s = slide()
title(s, "The change: behavioural feature engineering", "From 80 raw shares → ~6 interpretable scalars")
tb(s, 0.6, 1.2, 9.0, 0.55, [[
    ("Summarise each curve into contrast-type features that ", 13, BODY, False, False),
    ("cancel the shared envelope", 13, ACCENT, True, False),
    (" and keep only what distinguishes roles:", 13, BODY, False, False)]], line=1.05)
feats = [
    ("A_net", "net directionality", "(entry−exit)/total → origin vs destination"),
    ("F_flip", "direction flip", "early vs late directionality → 'leave then return' vs 'arrive then leave'"),
    ("t_median", "timing", "when activity peaks — early vs deep night"),
    ("persist_00", "late persistence", "share of activity after midnight"),
    ("hump_22", "nightlife bump", "secondary 22–23h peak"),
    ("dawn", "dawn rebound", "04:30–06:00 share (early-shift / airport)"),
]
x, y = 0.5, 1.95
for i,(k_,n_,d_) in enumerate(feats):
    col = i % 3; row = i // 3
    bx = 0.5 + col*3.05; by = 1.95 + row*1.35
    rect(s, bx, by, 2.9, 1.18, TINT)
    tb(s, bx+0.12, by+0.1, 2.66, 0.35, [[(k_, 12.5, ACCENT, True, False), ("  "+n_, 10.5, MUTED, False, True)]])
    tb(s, bx+0.12, by+0.5, 2.66, 0.65, [[(d_, 10, BODY, False, False)]], line=1.0)
tb(s, 0.6, 4.75, 9.0, 0.5, [[
    ("Contrast features (A_net, F_flip) subtract the common decline, isolating direction & timing — the axes that actually define night-time roles.", 11, MUTED, False, True)]], line=1.0)

# ---------- Slide 7: effect ----------
s = slide()
title(s, "Effect of the change", "Separability jumped; K became evidence-based")
# big stat
rect(s, 0.5, 1.3, 2.7, 1.5, ACCENT)
tb(s, 0.6, 1.45, 2.5, 0.9, [[("0.07 → 0.33", 26, WHITE, True, False)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
tb(s, 0.6, 2.3, 2.5, 0.4, [[("weekday silhouette\n(raw → v2 features)", 10.5, TINT, False, False)]], align=PP_ALIGN.CENTER)
for i,(t_) in enumerate([
    "Clusters now have distinct, nameable shapes",
    "GMM, KMeans & Ward broadly agree on the structure",
    "Bootstrap stability (ARI) used to choose K objectively",
    "Rail weekday K=4 stable (ARI 0.83); bus K=3 very stable (0.84–0.96)",
]):
    rect(s, 3.45, 1.3+i*0.55, 0.08, 0.42, ACC2)
    tb(s, 3.62, 1.28+i*0.55, 5.9, 0.55, [[(t_, 11.5, BODY, False, False)]], line=1.0)
pic_fit(s, P(r"lu_rail_2000_featv2\candidates\weekday_k4_profiles.png"), 3.4, 3.5, 6.1, 1.95,
        caption=None)
tb(s, 0.5, 3.6, 2.8, 1.6, [[("Same data, redesigned features → four clearly different weekday rail profiles (right).", 11, MUTED, False, True)]], line=1.05)

# ---------- Slide 7b: K-selection evidence ----------
s = slide()
title(s, "Choosing K — convergent evidence", "Silhouette + bootstrap stability + interpretability, per mode")
hdr = ["", "Chosen K", "Silhouette", "Stability (ARI)", "Why"]
rows = [
    ["Rail · weekday", "4", "0.33", "0.83", "4 distinct roles; K=3 dissolves nightlife"],
    ["Rail · weekend", "3 + airport", "0.23", "0.47", "weaker structure; Heathrow flagged separately"],
    ["Bus · weekday", "3", "0.20", "0.84", "even, stable; higher K only peels outliers"],
    ["Bus · weekend", "3", "0.22", "0.96", "very robust; mostly timing-driven"],
]
colx = [0.5, 2.7, 3.9, 5.2, 6.7]
colw = [2.2, 1.2, 1.3, 1.5, 3.3]
# header row
rect(s, 0.5, 1.35, 9.0, 0.45, ACCENT)
for j,h_ in enumerate(hdr):
    tb(s, colx[j]+0.05, 1.4, colw[j], 0.36, [[(h_, 11, WHITE, True, False)]], anchor=MSO_ANCHOR.MIDDLE)
for i,row in enumerate(rows):
    y = 1.8 + i*0.62
    if i % 2 == 0:
        rect(s, 0.5, y, 9.0, 0.62, TINT)
    for j,val in enumerate(row):
        bold = (j == 1)
        col = ACCENT if j == 1 else BODY
        sz = 12.5 if j == 1 else (11 if j != 4 else 10)
        tb(s, colx[j]+0.05, y+0.06, colw[j], 0.5, [[(val, sz, col, bold, False)]], anchor=MSO_ANCHOR.MIDDLE, line=1.0)
tb(s, 0.6, 4.5, 9.0, 0.7, [[
    ("Rule: ", 12, ACCENT, True, False),
    ("when indices are flat (continuous structure), stability + interpretability decide — not a silhouette peak. Bus silhouettes sit lower because LSOA areal data is noisier than point stations.", 12, BODY, False, False)]], line=1.05)

# ---------- Slide 8: rail results ----------
s = slide()
title(s, "Tentative results — Rail (LU)", "Weekday K=4: four interpretable night-time roles")
pic_fit(s, P(r"lu_rail_2000_featv2\candidates\weekday_k4_map.png"), 0.3, 1.2, 4.6, 4.2,
        caption="Weekday K=4 — spatial distribution")
roles = [
    ("C0 · Nightlife / destination", "central, zone 1; arrivals + 22h entry bump"),
    ("C1 · Residential return (early)", "strong exit-dominant, early flip; outer"),
    ("C2 · Balanced / through", "entry≈exit, smooth decline"),
    ("C3 · Residential return (later)", "exit-dominant, later crossover"),
]
tb(s, 5.0, 1.3, 4.7, 0.4, [[("Four roles", 14, ACCENT, True, False)]])
for i,(h_,b_) in enumerate(roles):
    y = 1.8 + i*0.72
    rect(s, 5.0, y, 0.08, 0.56, ACCENT)
    tb(s, 5.18, y-0.02, 4.5, 0.4, [[(h_, 12, ACCENT, True, False)]])
    tb(s, 5.18, y+0.32, 4.5, 0.35, [[(b_, 10.5, BODY, False, False)]], line=1.0)
tb(s, 5.0, 4.75, 4.7, 0.6, [[
    ("Weekend → K=3 main types + Heathrow handled separately (next slide).", 11, MUTED, False, True)]], line=1.05)

# ---------- Slide 9: Heathrow highlight ----------
s = slide()
title(s, "A validation moment — Heathrow", "The method cleanly isolates a real special type")
tb(s, 0.6, 1.25, 9.0, 0.8, [[
    ("At weekend K=5, the smallest cluster is exactly the three ", 14, BODY, False, False),
    ("Heathrow airport stations", 14, ACCENT, True, False),
    (" — not noise, but a genuine functional type with a unique signature.", 14, BODY, False, False)]], line=1.1)
rect(s, 0.5, 2.2, 9.0, 1.15, TINT)
tb(s, 0.7, 2.32, 8.6, 0.4, [[("Heathrow Terminals 1-2-3  ·  Terminal 4  ·  Terminal 5", 13, ACCENT, True, False)]])
tb(s, 0.7, 2.78, 8.6, 0.5, [[
    ("Extreme F_flip (0.6–1.15, highest in the network) + high dawn share — evening arrivals, strong pre-dawn departure surge (early flights / shift workers).", 11, BODY, False, False)]], line=1.05)
take = [
    ("Treat as a flagged special type", "Use K=3 main structure + label Heathrow explicitly as 'airport / 24h' rather than chasing it via higher K."),
    ("Evidence the features work", "A_net + F_flip + dawn pick out a real-world generator unsupervised — strong methodological validation."),
]
x = 0.5
for h_,b_ in take:
    rect(s, x, 3.55, 4.45, 1.5, WHITE, line=ACC2)
    tb(s, x+0.18, 3.68, 4.1, 0.5, [[(h_, 12.5, ACCENT, True, False)]])
    tb(s, x+0.18, 4.15, 4.1, 0.85, [[(b_, 11, BODY, False, False)]], line=1.05)
    x += 4.6

# ---------- Slide 9b: weekend results (both modes) ----------
s = slide()
title(s, "Weekend results — both modes", "Less differentiated than weekday → fewer types")
pic_fit(s, P(r"lu_rail_2000_featv2\candidates\weekend_k3_map.png"), 0.25, 1.25, 4.6, 3.55,
        caption="Rail weekend K=3 (+ Heathrow as special)")
pic_fit(s, P(r"busto_lsoa_2000_featv2\candidates\weekend_k3_map.png"), 5.05, 1.25, 4.7, 3.55,
        caption="Bus weekend K=3 (LSOA)")
tb(s, 0.6, 4.85, 9.0, 0.6, [[
    ("Weekend night travel is more homogeneous: ", 11.5, ACCENT, True, False),
    ("both modes settle at K=3; the weekday rail nightlife/residential split softens, and the dawn type carries the equity signal.", 11.5, BODY, False, False)]], line=1.05)

# ---------- Slide 10: bus results (v2 K=3) ----------
s = slide()
title(s, "Tentative results — Bus (LSOA)", "Weekday K=3 (v2 features) — three stable, balanced types")
pic_fit(s, P(r"busto_lsoa_2000_featv2\candidates\weekday_k3_map.png"), 0.25, 1.2, 4.8, 4.25,
        caption="Weekday K=3 — LSOA cluster map (v2 features)")
tb(s, 5.15, 1.3, 4.6, 0.4, [[("Three bus types (K=3)", 13.5, ACCENT, True, False)]])
for i,(h_,b_) in enumerate([
    ("Destination-early", "alighting-dominant, fades early"),
    ("Balanced", "boardings ≈ alightings"),
    ("Late / dawn-persistent", "arrival + strong 05:00 boarding surge"),
]):
    y = 1.75 + i*0.62
    rect(s, 5.15, y, 0.08, 0.48, ACCENT)
    tb(s, 5.33, y-0.02, 4.4, 0.35, [[(h_, 11.5, ACCENT, True, False)]])
    tb(s, 5.33, y+0.27, 4.4, 0.3, [[(b_, 10, BODY, False, False)]], line=1.0)
rect(s, 5.15, 3.75, 4.55, 1.35, TINT)
tb(s, 5.33, 3.86, 4.2, 0.4, [[("Balanced & robust", 12, ACCENT, True, False)]])
tb(s, 5.33, 4.24, 4.2, 0.85, [[
    ("Cluster sizes even (1180 / 975 / 831) and bootstrap-stable (ARI 0.84) — no degenerate splits, unlike forcing higher K.", 10.5, BODY, False, False)]], line=1.05)

# ---------- Slide 10b: bus night-service coverage (equity) ----------
s = slide()
title(s, "Bus night-service coverage — equity read-out", "Before clustering: where is there any night bus at all?")
pic_fit(s, P(r"busto_lsoa_2000_gmm\weekday\weekday_bus_cluster_map.png"), 0.25, 1.2, 4.9, 4.25,
        caption="Weekday: clustered / low-service / no-service LSOAs")
tb(s, 5.25, 1.35, 4.5, 0.4, [[("4,994 London LSOAs split three ways", 13, ACCENT, True, False)]])
stats = [("60%", "clustered", "≥ service threshold (coloured)"),
         ("22%", "low service", "some night bus, below threshold"),
         ("18%", "no service", "no night bus recorded at all")]
for i,(n_,h_,b_) in enumerate(stats):
    y = 1.9 + i*0.78
    rect(s, 5.25, y, 1.0, 0.66, ACCENT)
    tb(s, 5.25, y+0.04, 1.0, 0.58, [[(n_, 19, WHITE, True, False)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    tb(s, 6.4, y-0.02, 3.3, 0.35, [[(h_, 12.5, ACCENT, True, False)]])
    tb(s, 6.4, y+0.32, 3.3, 0.35, [[(b_, 10.5, BODY, False, False)]], line=1.0)
tb(s, 5.25, 4.4, 4.5, 0.7, [[
    ("Low/no-service LSOAs cluster in the outer ring — a direct night-time accessibility gap, and an input to RQ3.", 11, MUTED, False, True)]], line=1.05)

# ---------- Slide 11: rail vs bus ----------
s = slide()
title(s, "Why rail splits more than bus", "Point stations vs aggregated areas (MAUP)")
cols = [
    ("Rail — discrete stations", ACCENT, [
        "Each station = one functional entity",
        "Sharp, distinctive signatures",
        "Airport / CBD / nightlife / residential",
        "Supports up to 4 types + special cases",
    ]),
    ("Bus — LSOA aggregates", ACC2, [
        "Each LSOA = many stops, mixed routes",
        "Distinctive signals averaged out",
        "More homogeneous / continuous",
        "Supports ~3 types; stays low-K",
    ]),
]
x = 0.5
for h_,c_,items in cols:
    rect(s, x, 1.3, 4.45, 3.4, TINT)
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(1.3), Inches(4.45), Inches(0.5))
    bar.fill.solid(); bar.fill.fore_color.rgb = c_; bar.line.fill.background(); bar.shadow.inherit=False
    tfb=bar.text_frame; pb=tfb.paragraphs[0]; pb.alignment=PP_ALIGN.CENTER
    rr=pb.add_run(); rr.text=h_; rr.font.name="Calibri"; rr.font.size=Pt(13.5); rr.font.bold=True; rr.font.color.rgb=WHITE
    for i,it in enumerate(items):
        tb(s, x+0.25, 2.0+i*0.62, 4.0, 0.55, [[("• ", 12, c_, True, False),(it, 11.5, BODY, False, False)]], line=1.0)
    x += 4.6
tb(s, 0.6, 4.85, 9.0, 0.5, [[
    ("This asymmetry is itself a finding: rail carries London's specialised night roles; bus is the more uniform spatial baseline.", 11, MUTED, False, True)]], line=1.0)

# ---------- Slide 12: status & next ----------
s = slide(DARK)
accent_bar2 = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.4), Inches(0.30), Inches(0.08), Inches(0.42))
accent_bar2.fill.solid(); accent_bar2.fill.fore_color.rgb = ACC2; accent_bar2.line.fill.background(); accent_bar2.shadow.inherit=False
tb(s, 0.6, 0.22, 9.0, 0.6, [[("Current status & next steps", 23, WHITE, True, False)]])
tb(s, 0.6, 0.95, 9.0, 0.4, [[("Tentative K (evidence-based)", 12.5, ACC2, False, True)]])
# K table as cards
ktab = [("Rail","weekday","K = 4"),("Rail","weekend","K = 3 + airport"),("Bus","weekday","K = 3"),("Bus","weekend","K = 3")]
x=0.5
for m_,d_,k_ in ktab:
    rect(s, x, 1.45, 2.2, 1.0, RGBColor(0x2E,0x23,0x55))
    tb(s, x+0.12, 1.55, 1.96, 0.3, [[(m_+" · "+d_, 10.5, ACC2, False, False)]])
    tb(s, x+0.12, 1.9, 1.96, 0.4, [[(k_, 15, WHITE, True, False)]])
    x+=2.37
tb(s, 0.6, 2.75, 9.0, 0.4, [[("Next", 13, ACC2, True, False)]])
for i,t_ in enumerate([
    "Finalise the four label sets + name every cluster; flag Heathrow as a special airport type",
    "Quantify spatial structure: Moran's I / LISA (bus areas), join-count (rail points)",
    "Link clusters to night-worker classification & socio-economic context (RQ2)",
    "Carry forward to demand–service mismatch analysis (RQ3)",
]):
    rect(s, 0.6, 3.2+i*0.52, 0.08, 0.4, ACC2)
    tb(s, 0.78, 3.18+i*0.52, 8.8, 0.5, [[(t_, 12, WHITE, False, False)]], line=1.0)

prs.save(OUT)
print("Saved", OUT, "slides:", len(prs.slides.__iter__.__self__._sldIdLst))
