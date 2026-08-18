# -*- coding: utf-8 -*-
"""Discussion-oriented deck: methods & results, with open questions for supervisor.
Overwrites night_transport_RQ1_progress.pptx."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image
import os

OUT = r"D:\SDS2025_workspace\CASA_FYP\FYP\night_transport_RQ1_progress.pptx"
IMG = r"D:\SDS2025_workspace\CASA_FYP\FYP\outputs"

DARK   = RGBColor(0x1F, 0x16, 0x40)
ACCENT = RGBColor(0x50, 0x07, 0x78)
ACC2   = RGBColor(0x8E, 0x6B, 0xB8)
BODY   = RGBColor(0x33, 0x33, 0x33)
MUTED  = RGBColor(0x66, 0x66, 0x66)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
TINT   = RGBColor(0xF4, 0xF0, 0xF8)
QFILL  = RGBColor(0xFB, 0xF3, 0xDE)   # discussion callout (warm)
QEDGE  = RGBColor(0xC8, 0x9B, 0x3C)
QTXT   = RGBColor(0x7A, 0x5A, 0x12)

prs = Presentation()
prs.slide_width = Inches(10); prs.slide_height = Inches(5.625)
BLANK = prs.slide_layouts[6]


def slide(bg=WHITE):
    s = prs.slides.add_slide(BLANK)
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    r.fill.solid(); r.fill.fore_color.rgb = bg; r.line.fill.background(); r.shadow.inherit = False
    s.shapes._spTree.remove(r._element); s.shapes._spTree.insert(2, r._element)
    return s


def tb(s, l, t, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, sp_after=4, line=1.0):
    box = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = box.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = Inches(0.05); tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02); tf.margin_bottom = Inches(0.02)
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.space_after = Pt(sp_after); p.space_before = Pt(0); p.line_spacing = line
        for (text, size, color, bold, italic) in para:
            r = p.add_run(); r.text = text
            r.font.name = "Calibri"; r.font.size = Pt(size); r.font.color.rgb = color
            r.font.bold = bold; r.font.italic = italic
    return box


def rect(s, l, t, w, h, fill, line=None, rounded=True):
    shp = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    r = s.shapes.add_shape(shp, Inches(l), Inches(t), Inches(w), Inches(h))
    r.fill.solid(); r.fill.fore_color.rgb = fill
    if line is None:
        r.line.fill.background()
    else:
        r.line.color.rgb = line; r.line.width = Pt(1.25)
    r.shadow.inherit = False
    return r


def title(s, text, sub=None):
    b = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.4), Inches(0.30), Inches(0.08), Inches(0.42))
    b.fill.solid(); b.fill.fore_color.rgb = ACCENT; b.line.fill.background(); b.shadow.inherit = False
    tb(s, 0.6, 0.22, 9.0, 0.6, [[(text, 22, ACCENT, True, False)]])
    if sub:
        tb(s, 0.6, 0.80, 9.0, 0.4, [[(sub, 12.5, MUTED, False, True)]])


def discuss(s, l, t, w, h, q):
    rect(s, l, t, w, h, QFILL, line=QEDGE)
    tb(s, l+0.18, t+0.08, w-0.35, 0.3, [[("◆  For discussion", 11, QTXT, True, False)]])
    tb(s, l+0.18, t+0.42, w-0.35, h-0.5, [[(q, 11.5, BODY, False, False)]], line=1.05)


def pic_fit(s, path, l, t, maxw, maxh, caption=None):
    iw, ih = Image.open(path).size
    ar = iw/ih; w, h = maxw, maxw/ar
    if h > maxh:
        h = maxh; w = maxh*ar
    left = l + (maxw-w)/2
    s.shapes.add_picture(path, Inches(left), Inches(t), Inches(w), Inches(h))
    if caption:
        tb(s, l, t+h+0.02, maxw, 0.3, [[(caption, 8.5, MUTED, False, True)]], align=PP_ALIGN.CENTER)
    return h


def P(sub):
    return os.path.join(IMG, sub)


# ===== 1. Title =====
s = slide(DARK)
tb(s, 0.6, 1.5, 8.8, 0.6, [[("Night-time Public Transport & Spatial Equity in London", 26, WHITE, True, False)]])
tb(s, 0.6, 2.5, 8.8, 0.5, [[("RQ1 — Methods & Results: points for discussion", 15, ACC2, True, False)]])
tb(s, 0.6, 3.25, 8.8, 0.9, [[("Window choice · clustering method · feature engineering · K & spatial unit", 12.5, RGBColor(0xCF,0xC4,0xE4), False, True)]])
tb(s, 0.6, 4.6, 8.8, 0.4, [[("Fangzheng  |  UCL CASA × TfL", 12, RGBColor(0xB9,0xAD,0xD6), False, False)]])

# ===== 2. Window change =====
s = slide()
title(s, "Decision 1 — analysis window: 18:00–06:00 → 20:00–06:00", "Why I narrowed the start of the night window")
# timeline bar
y = 1.55
rect(s, 0.6, y, 2.4, 0.5, RGBColor(0xDD,0xDD,0xDD))           # 18-20 excluded
rect(s, 3.0, y, 6.4, 0.5, ACCENT)                              # 20-06 kept
tb(s, 0.6, y+0.08, 2.4, 0.35, [[("18:00–20:00  excluded", 10.5, MUTED, True, False)]], align=PP_ALIGN.CENTER)
tb(s, 3.0, y+0.08, 6.4, 0.35, [[("20:00–06:00  analysed", 11.5, WHITE, True, False)]], align=PP_ALIGN.CENTER)
tb(s, 0.6, y+0.52, 9.0, 0.3, [[("18:00      20:00            23:00            00:00            03:00            06:00", 9, MUTED, False, False)]])
pts = [
    ("18:00–20:00 is still the PM commute tail", "a strong, generic commute signal that dominates the clustering and masks night-specific patterns"),
    ("20:00 better matches the 'night-time economy'", "leisure / night-work travel rather than the end of the working day"),
    ("Keeps the window comparable across modes & days", "same 20:00 start for rail and bus, weekday and weekend"),
]
yy = 2.35
for h_, b_ in pts:
    rect(s, 0.6, yy, 0.08, 0.55, ACCENT)
    tb(s, 0.8, yy-0.02, 8.6, 0.3, [[(h_, 12.5, ACCENT, True, False)]])
    tb(s, 0.8, yy+0.30, 8.6, 0.3, [[(b_, 11, BODY, False, False)]], line=1.0)
    yy += 0.66
discuss(s, 0.6, 4.45, 8.8, 0.85,
        "Is 20:00 the right cut-off, or should the night window be defined differently (e.g. 22:00 / 23:00, or aligned to Night Tube hours)? Does excluding 18:00–20:00 lose anything important?")

# ===== 3. Clustering method (GMM-full) =====
s = slide()
title(s, "Decision 2 — clustering algorithm", "GMM-full underperformed early on; what should we standardise on?")
tb(s, 0.6, 1.2, 9.0, 0.7, [[
    ("Early attempts used ", 13, BODY, False, False),
    ("GMM with full covariance", 13, ACCENT, True, False),
    (" on the high-dimensional raw features — it was unstable (near-singular covariances, parameters ≫ 270 stations).", 13, BODY, False, False)]], line=1.1)
cols = [
    ("Then", "Switched to GMM (diagonal cov) for stability — but K stayed ambiguous on the raw features."),
    ("Now", "With low-dim engineered features, the issue is largely resolved: GMM, KMeans & Ward broadly agree (silhouette ≈ 0.26–0.34)."),
]
x = 0.6
for h_, b_ in cols:
    rect(s, x, 2.05, 4.35, 1.25, TINT)
    tb(s, x+0.18, 2.16, 4.0, 0.35, [[(h_, 13, ACCENT, True, False)]])
    tb(s, x+0.18, 2.55, 4.0, 0.7, [[(b_, 11, BODY, False, False)]], line=1.05)
    x += 4.55
tb(s, 0.6, 3.5, 9.0, 0.4, [[("Cross-model silhouette on v2 features (rail weekday): GMM 0.26–0.32 · KMeans 0.33–0.34 · Ward 0.32", 10.5, MUTED, False, True)]])
discuss(s, 0.6, 4.0, 8.8, 1.25,
        "Which method do you recommend we commit to — model-based GMM (soft assignment + BIC) or KMeans/Ward (simpler, here slightly higher silhouette)? Is reporting one primary + one robustness check acceptable?")

# ===== 4. Feature engineering — motivation =====
s = slide()
title(s, "Decision 3 — feature engineering (motivation)", "Early version clustered on raw normalised entry/exit shares only")
pic_fit(s, P(r"lu_rail_2000_diag\weekday\weekday_profiles.png"), 0.35, 1.25, 4.7, 4.0,
        caption="Old features (raw shares): clusters nearly identical")
tb(s, 5.3, 1.3, 4.4, 0.4, [[("The problem", 14, ACCENT, True, False)]])
for i, t_ in enumerate([
    "Every curve shares the same night-decline shape ('envelope')",
    "That common shape is the biggest source of variance → dominates distance",
    "Discriminative signals (direction, timing) get buried",
    "Result: near-duplicate clusters, silhouette ≈ 0.07, K unresolvable",
]):
    rect(s, 5.3, 1.85+i*0.66, 0.08, 0.5, ACCENT)
    tb(s, 5.48, 1.83+i*0.66, 4.0, 0.62, [[(t_, 11, BODY, False, False)]], line=1.0)
tb(s, 5.3, 4.55, 4.4, 0.6, [[
    ("→ After consulting on approach, moved to ", 11.5, BODY, False, True),
    ("behavioural feature engineering", 11.5, ACCENT, True, True),
    (".", 11.5, BODY, False, True)]], line=1.05)

# ===== 5. Feature engineering — list & selection =====
s = slide()
title(s, "Feature engineering — feature list & selection logic", "Summarise each curve into a few contrast-type scalars")
feats = [
    ("A_net", "net directionality", "(entry−exit)/total → origin vs destination"),
    ("F_flip", "direction flip", "early vs late directionality → leave-then-return vs arrive-then-leave"),
    ("t_median", "timing", "when activity peaks (early vs deep night)"),
    ("persist_00", "late persistence", "share of activity after midnight"),
    ("hump_22", "nightlife bump", "secondary 22–23h peak"),
    ("dawn", "dawn rebound", "04:30–06:00 share (early-shift / airport)"),
]
for i, (k_, n_, d_) in enumerate(feats):
    col = i % 3; row = i // 3
    bx = 0.55 + col*3.05; by = 1.2 + row*1.05
    rect(s, bx, by, 2.9, 0.92, TINT)
    tb(s, bx+0.12, by+0.07, 2.66, 0.3, [[(k_, 12, ACCENT, True, False), ("  "+n_, 9.5, MUTED, False, True)]])
    tb(s, bx+0.12, by+0.40, 2.66, 0.5, [[(d_, 9.5, BODY, False, False)]], line=0.95)
tb(s, 0.55, 3.4, 9.0, 0.35, [[("Selection logic", 13, ACCENT, True, False)]])
for i, t_ in enumerate([
    "Prefer contrast features (A_net, F_flip): subtracting entry/exit cancels the shared envelope",
    "Drop dead / redundant axes: 'dawn' is empty on weekday; 'early_conc' ≈ −t_median (corr .85–.91), pruned",
    "Standardise (z-score) so axes weigh equally; check correlation matrix before clustering",
]):
    rect(s, 0.55, 3.78+i*0.5, 0.08, 0.38, ACC2)
    tb(s, 0.73, 3.76+i*0.5, 8.9, 0.45, [[(t_, 11, BODY, False, False)]], line=1.0)

# ===== 6. Feature engineering — A_net & F_flip detail =====
s = slide()
title(s, "The two core axes — A_net & F_flip", "These carry most of the night-time 'role' signal")
rect(s, 0.55, 1.25, 4.4, 2.0, TINT)
tb(s, 0.73, 1.35, 4.05, 0.35, [[("A_net — net directionality", 13, ACCENT, True, False)]])
tb(s, 0.73, 1.75, 4.05, 0.5, [[("A_net = (Σentry − Σexit) / (Σentry + Σexit)", 11.5, BODY, True, False)]])
tb(s, 0.73, 2.2, 4.05, 1.0, [[
    (">0 entry-dominant = origin (people leaving);  <0 exit-dominant = destination (people arriving). "
     "Dividing by total makes it scale-free and cancels the shared decline.", 10.5, BODY, False, False)]], line=1.05)
rect(s, 5.05, 1.25, 4.4, 2.0, TINT)
tb(s, 5.23, 1.35, 4.05, 0.35, [[("F_flip — direction flip", 13, ACCENT, True, False)]])
tb(s, 5.23, 1.75, 4.05, 0.5, [[("F_flip = A_early − A_late   (early=20–23h, late=23–01h)", 11, BODY, True, False)]])
tb(s, 5.23, 2.2, 4.05, 1.0, [[
    (">0 early-out → late-in ('leave then return', residential);  <0 early-in → late-out "
     "('arrive then leave', nightlife/destination). Signed version of the entry/exit crossover.", 10.5, BODY, False, False)]], line=1.05)
tb(s, 0.55, 3.4, 9.0, 0.5, [[
    ("Together they place each station on two interpretable axes; e.g. Heathrow sits at the extreme of F_flip (airport).", 11, MUTED, False, True)]])
discuss(s, 0.55, 3.95, 8.9, 1.3,
        "Are these engineered features methodologically sound, or do they over-inject prior assumptions? "
        "Which would you add / drop? Should raw-share + PCA be reported alongside as a more 'hands-off' check, "
        "and does compositional (log-ratio) treatment matter here?")

# ===== 7. Effect =====
s = slide()
title(s, "Effect of the change", "Separability jumped; a real special type emerged")
rect(s, 0.55, 1.3, 2.6, 1.35, ACCENT)
tb(s, 0.6, 1.45, 2.5, 0.8, [[("0.07 → 0.33", 24, WHITE, True, False)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
tb(s, 0.6, 2.25, 2.5, 0.35, [[("weekday silhouette", 10, TINT, False, False)]], align=PP_ALIGN.CENTER)
for i, t_ in enumerate([
    "Distinct, nameable cluster shapes",
    "GMM / KMeans / Ward broadly agree",
    "Heathrow's 3 stations isolated cleanly\n(extreme F_flip + dawn) — feature validation",
]):
    rect(s, 3.35, 1.3+i*0.46, 0.08, 0.36, ACC2)
    tb(s, 3.52, 1.27+i*0.46, 6.0, 0.5, [[(t_, 11, BODY, False, False)]], line=1.0)
pic_fit(s, P(r"lu_rail_2000_featv2\candidates\weekday_k4_profiles.png"), 3.3, 2.95, 6.3, 2.4,
        caption="Rail weekday K=4 — four clearly different profiles")
tb(s, 0.55, 3.0, 2.65, 2.0, [[("Same data, redesigned features.", 11, MUTED, False, True)]], line=1.05)

# ===== 8. K selection =====
s = slide()
title(s, "Open question — choosing K", "Indices are flat (continuous structure); how to decide?")
hdr = ["", "Tent. K", "Silhouette", "Stability", "Note"]
rows = [
    ["Rail · weekday", "4", "0.33", "0.83", "4 distinct roles; K=3 dissolves nightlife"],
    ["Rail · weekend", "3 +airport", "0.23", "0.47", "weaker; Heathrow flagged separately"],
    ["Bus · weekday", "3", "0.20", "0.84", "even, stable; higher K peels outliers"],
    ["Bus · weekend", "3", "0.22", "0.96", "very robust; timing-driven"],
]
colx = [0.55, 2.5, 3.7, 5.0, 6.2]; colw = [2.0, 1.2, 1.3, 1.2, 3.3]
rect(s, 0.55, 1.25, 8.95, 0.42, ACCENT)
for j, h_ in enumerate(hdr):
    tb(s, colx[j]+0.05, 1.29, colw[j], 0.34, [[(h_, 10.5, WHITE, True, False)]], anchor=MSO_ANCHOR.MIDDLE)
for i, row in enumerate(rows):
    yv = 1.67 + i*0.46
    if i % 2 == 0:
        rect(s, 0.55, yv, 8.95, 0.46, TINT)
    for j, val in enumerate(row):
        bold = (j == 1); col = ACCENT if j == 1 else BODY
        sz = 12 if j == 1 else (10.5 if j != 4 else 9.5)
        tb(s, colx[j]+0.05, yv+0.04, colw[j], 0.4, [[(val, sz, col, bold, False)]], anchor=MSO_ANCHOR.MIDDLE, line=1.0)
discuss(s, 0.55, 3.7, 8.95, 1.5,
        "When silhouette is flat, we lean on bootstrap stability + interpretability rather than an index peak. "
        "Is that defensible? Are the tentative Ks (rail 4 / weekend 3 + airport; bus 3) reasonable, or would you "
        "expect a different number of night-time types?")

# ===== 9. Bus spatial unit =====
s = slide()
title(s, "Open question — is LSOA-level bus clustering sound?", "Point stations vs aggregated areas (MAUP)")
colsd = [
    ("Rail — discrete stations", ACCENT, ["Each station = one function", "Sharp signatures", "Airport / CBD / nightlife / residential", "Up to 4 types + specials"]),
    ("Bus — LSOA aggregates", ACC2, ["Each LSOA = many stops, mixed routes", "Distinctive signals averaged out", "More homogeneous / continuous", "~3 types; stays low-K"]),
]
x = 0.55
for h_, c_, items in colsd:
    rect(s, x, 1.25, 4.35, 2.3, TINT)
    bar = rect(s, x, 1.25, 4.35, 0.45, c_, rounded=False)
    tf = bar.text_frame; p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    rr = p.add_run(); rr.text = h_; rr.font.name = "Calibri"; rr.font.size = Pt(12.5); rr.font.bold = True; rr.font.color.rgb = WHITE
    for i, it in enumerate(items):
        tb(s, x+0.22, 1.85+i*0.42, 3.95, 0.4, [[("• ", 11, c_, True, False), (it, 10.5, BODY, False, False)]], line=1.0)
    x += 4.55
discuss(s, 0.55, 3.75, 8.95, 1.45,
        "Is LSOA the right unit for bus, or should we cluster at stop level (finer, but noisier) or another aggregation? "
        "The point/area asymmetry means bus shows fewer types than rail — is that an acceptable finding or a unit artefact to fix?")

# ===== 10. Bus blank areas =====
s = slide()
title(s, "Open question — bus 'blank' areas (n ≥ 50 filter)", "40% of LSOAs drop out before clustering")
pic_fit(s, P(r"busto_lsoa_2000_gmm\weekday\weekday_bus_cluster_map.png"), 0.25, 1.2, 4.7, 4.2,
        caption="Weekday: clustered / low-service / no-service LSOAs")
tb(s, 5.05, 1.3, 4.6, 0.4, [[("4,994 LSOAs → split three ways", 13, ACCENT, True, False)]])
for i, (n_, h_) in enumerate([("60%", "clustered (≥50 night demand)"), ("22%", "low service (0 < demand < 50)"), ("18%", "no night bus recorded")]):
    yv = 1.8 + i*0.6
    rect(s, 5.05, yv, 0.95, 0.5, ACCENT)
    tb(s, 5.05, yv+0.03, 0.95, 0.44, [[(n_, 16, WHITE, True, False)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    tb(s, 6.15, yv+0.06, 3.5, 0.4, [[(h_, 11, BODY, False, False)]], anchor=MSO_ANCHOR.MIDDLE, line=1.0)
discuss(s, 5.05, 3.65, 4.55, 1.6,
        "How should the dropped 40% be handled? Options: keep them as explicit 'low/no-service' categories "
        "(an equity finding), lower the n≥50 threshold, or change the spatial unit. "
        "Is the threshold itself defensible?")

# ===== 11. Current tentative results =====
s = slide()
title(s, "Where the results stand right now", "Tentative, pending the decisions above")
pic_fit(s, P(r"lu_rail_2000_featv2\candidates\weekday_k4_map.png"), 0.3, 1.25, 4.5, 3.5,
        caption="Rail weekday K=4")
pic_fit(s, P(r"busto_lsoa_2000_featv2\candidates\weekday_k3_map.png"), 5.0, 1.25, 4.6, 3.5,
        caption="Bus weekday K=3")
tb(s, 0.6, 4.8, 9.0, 0.6, [[
    ("Rail: nightlife / residential-early / balanced / residential-late (+ airport).   "
     "Bus: destination-early / balanced / late-dawn.", 11, BODY, False, False)]], line=1.05)

# ===== 12. Discussion summary =====
s = slide(DARK)
b = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.4), Inches(0.30), Inches(0.08), Inches(0.42))
b.fill.solid(); b.fill.fore_color.rgb = ACC2; b.line.fill.background(); b.shadow.inherit = False
tb(s, 0.6, 0.22, 9.0, 0.6, [[("Questions for today", 22, WHITE, True, False)]])
qs = [
    ("Window", "Is 20:00–06:00 the right night window?"),
    ("Method", "Which clustering algorithm should we standardise on?"),
    ("Features", "Is the behavioural feature engineering sound — what to add / drop?"),
    ("K", "Are the tentative Ks (rail 4 / bus 3) reasonable when indices are flat?"),
    ("Spatial unit", "Is LSOA-level bus clustering appropriate?"),
    ("Blank areas", "How to handle the 40% of LSOAs dropped by the n≥50 filter?"),
]
for i, (h_, q_) in enumerate(qs):
    row = i // 2; col = i % 2
    bx = 0.6 + col*4.6; by = 1.25 + row*1.25
    rect(s, bx, by, 4.4, 1.05, RGBColor(0x2E,0x23,0x55))
    tb(s, bx+0.18, by+0.1, 4.05, 0.35, [[(f"{i+1}. {h_}", 12.5, ACC2, True, False)]])
    tb(s, bx+0.18, by+0.45, 4.05, 0.55, [[(q_, 11, WHITE, False, False)]], line=1.0)

prs.save(OUT)
print("Saved", OUT, "slides:", len(prs.slides._sldIdLst))
